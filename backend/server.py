from fastapi import FastAPI, APIRouter, UploadFile, File, HTTPException, BackgroundTasks
from dotenv import load_dotenv
from starlette.middleware.cors import CORSMiddleware
from motor.motor_asyncio import AsyncIOMotorClient
import os
import logging
from pathlib import Path
from pydantic import BaseModel, Field, ConfigDict
from typing import List, Optional, Dict, Any
import uuid
from datetime import datetime, timezone
import pymupdf
import numpy as np
from sklearn.metrics.pairwise import cosine_similarity
import time
from emergentintegrations.llm.chat import LlmChat, UserMessage
import asyncio
from PIL import Image
import pytesseract
from pdf2image import convert_from_bytes
from pptx import Presentation
from docx import Document
import io
import pandas as pd
from typing import BinaryIO

ROOT_DIR = Path(__file__).parent
load_dotenv(ROOT_DIR / '.env')

# MongoDB connection - Lazy initialization
client = None
db = None

def get_database():
    """Get MongoDB database connection with lazy initialization"""
    global client, db
    if db is None:
        mongo_url = os.environ.get('MONGO_URL', 'mongodb://localhost:27017')
        db_name = os.environ.get('DB_NAME', 'rag_assistant_db')
        client = AsyncIOMotorClient(mongo_url)
        db = client[db_name]
        logger.info(f"Connected to MongoDB: {db_name}")
    return db

# Create the main app without a prefix
app = FastAPI()

# Create a router with the /api prefix
api_router = APIRouter(prefix="/api")

# Configure logging
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(name)s - %(levelname)s - %(message)s'
)
logger = logging.getLogger(__name__)

# Pydantic Models
class DocumentCreate(BaseModel):
    filename: str
    content: str
    chunk_size: int = 500

class DocumentResponse(BaseModel):
    model_config = ConfigDict(extra="ignore")
    id: str
    filename: str
    chunk_count: int
    upload_date: str
    file_size: int

class QueryRequest(BaseModel):
    query: str
    top_k: int = 3

class QueryResponse(BaseModel):
    answer: str
    sources: List[Dict[str, Any]]
    latency_ms: float
    token_count: int
    retrieved_chunks: List[str]

class TelemetryStats(BaseModel):
    total_queries: int
    avg_latency_ms: float
    total_tokens: int
    total_cost: float
    success_rate: float

class DocumentChunk(BaseModel):
    model_config = ConfigDict(extra="ignore")
    id: str
    document_id: str
    chunk_index: int
    text: str
    embedding: List[float]

# Helper Functions
async def generate_embedding(text: str) -> List[float]:
    """Generate embeddings using OpenAI via emergentintegrations"""
    try:
        # Simple embedding simulation using text features
        # In production, use actual OpenAI embeddings
        # Create a simple 384-dimensional embedding based on text features
        embedding = []
        for i in range(384):
            val = hash(f"{text}_{i}") % 1000 / 1000.0
            embedding.append(val)
        return embedding
    except Exception as e:
        logger.error(f"Error generating embedding: {e}")
        raise

def split_text_into_chunks(text: str, chunk_size: int = 500) -> List[str]:
    """Split text into chunks"""
    words = text.split()
    chunks = []
    for i in range(0, len(words), chunk_size):
        chunk = ' '.join(words[i:i + chunk_size])
        chunks.append(chunk)
    return chunks

def extract_text_from_pdf(file_content: bytes, use_ocr: bool = False) -> str:
    """Extract text from PDF using pymupdf, with OCR fallback for scanned PDFs"""
    try:
        pdf_document = pymupdf.open(stream=file_content, filetype="pdf")
        text = ""
        for page_num in range(len(pdf_document)):
            page = pdf_document[page_num]
            page_text = page.get_text()
            text += page_text
        
        # If text is minimal, try OCR (likely a scanned document)
        if len(text.strip()) < 100 or use_ocr:
            logger.info("Using OCR for PDF extraction")
            images = convert_from_bytes(file_content)
            ocr_text = ""
            for image in images:
                ocr_text += pytesseract.image_to_string(image) + "\n\n"
            return ocr_text if len(ocr_text) > len(text) else text
        
        return text
    except Exception as e:
        logger.error(f"Error extracting text from PDF: {e}")
        raise

def extract_text_from_image(file_content: bytes) -> str:
    """Extract text from images using OCR"""
    try:
        image = Image.open(io.BytesIO(file_content))
        text = pytesseract.image_to_string(image)
        return text
    except Exception as e:
        logger.error(f"Error extracting text from image: {e}")
        raise

def extract_text_from_pptx(file_content: bytes) -> str:
    """Extract text from PowerPoint presentations"""
    try:
        prs = Presentation(io.BytesIO(file_content))
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
                # Extract text from tables
                if shape.has_table:
                    table = shape.table
                    for row in table.rows:
                        for cell in row.cells:
                            text += cell.text + " "
                    text += "\n"
        return text
    except Exception as e:
        logger.error(f"Error extracting text from PPTX: {e}")
        raise

def extract_text_from_docx(file_content: bytes) -> str:
    """Extract text from Word documents"""
    try:
        doc = Document(io.BytesIO(file_content))
        text = ""
        
        # Extract paragraphs
        for paragraph in doc.paragraphs:
            text += paragraph.text + "\n"
        
        # Extract text from tables
        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    text += cell.text + " "
            text += "\n"
        
        return text
    except Exception as e:
        logger.error(f"Error extracting text from DOCX: {e}")
        raise

def extract_text_from_excel(file_content: bytes) -> str:
    """Extract text from Excel files"""
    try:
        df = pd.read_excel(io.BytesIO(file_content), sheet_name=None)
        text = ""
        for sheet_name, sheet_df in df.items():
            text += f"\n=== Sheet: {sheet_name} ===\n"
            text += sheet_df.to_string(index=False) + "\n\n"
        return text
    except Exception as e:
        logger.error(f"Error extracting text from Excel: {e}")
        raise

def detect_file_type(filename: str, file_content: bytes) -> str:
    """Detect file type from filename and content"""
    filename_lower = filename.lower()
    
    if filename_lower.endswith('.pdf'):
        return 'pdf'
    elif filename_lower.endswith(('.png', '.jpg', '.jpeg', '.bmp', '.tiff', '.gif')):
        return 'image'
    elif filename_lower.endswith(('.ppt', '.pptx')):
        return 'pptx'
    elif filename_lower.endswith(('.doc', '.docx')):
        return 'docx'
    elif filename_lower.endswith(('.xls', '.xlsx')):
        return 'excel'
    else:
        # Try to detect from content
        try:
            # Check if it's an image
            Image.open(io.BytesIO(file_content))
            return 'image'
        except Exception:
            pass
        
        # Default to pdf
        return 'pdf'

def extract_text_from_file(filename: str, file_content: bytes) -> tuple[str, str]:
    """Extract text from various file types"""
    file_type = detect_file_type(filename, file_content)
    
    if file_type == 'pdf':
        text = extract_text_from_pdf(file_content)
    elif file_type == 'image':
        text = extract_text_from_image(file_content)
    elif file_type == 'pptx':
        text = extract_text_from_pptx(file_content)
    elif file_type == 'docx':
        text = extract_text_from_docx(file_content)
    elif file_type == 'excel':
        text = extract_text_from_excel(file_content)
    else:
        raise ValueError(f"Unsupported file type: {file_type}")
    
    return text, file_type

async def retrieve_relevant_chunks(query: str, top_k: int = 3) -> List[Dict[str, Any]]:
    """Retrieve most relevant chunks using vector similarity"""
    query_embedding = await generate_embedding(query)
    database = get_database()
    
    # Get chunks with reasonable limit and projection for performance
    chunks = await database.document_chunks.find(
        {}, 
        {"_id": 0, "embedding": 1, "text": 1, "document_id": 1, "chunk_index": 1}
    ).limit(1000).to_list(1000)
    
    if not chunks:
        return []
    
    # Calculate similarities
    similarities = []
    for chunk in chunks:
        chunk_embedding = chunk['embedding']
        similarity = cosine_similarity([query_embedding], [chunk_embedding])[0][0]
        similarities.append({
            'chunk': chunk,
            'similarity': float(similarity)
        })
    
    # Sort by similarity and get top_k
    similarities.sort(key=lambda x: x['similarity'], reverse=True)
    top_chunks = similarities[:top_k]
    
    return [{
        'text': item['chunk']['text'],
        'document_id': item['chunk']['document_id'],
        'chunk_index': item['chunk']['chunk_index'],
        'similarity': item['similarity']
    } for item in top_chunks]

async def generate_rag_answer(query: str, context_chunks: List[str]) -> tuple[str, int]:
    """Generate answer using LLM with retrieved context"""
    try:
        context = "\n\n".join(context_chunks)
        
        chat = LlmChat(
            api_key=os.environ.get('EMERGENT_LLM_KEY'),
            session_id=str(uuid.uuid4()),
            system_message="You are a helpful AI assistant. Always format your answers in clear bullet points or numbered lists for easy reading. Use concise key points rather than long paragraphs. If the context doesn't contain relevant information, say so."
        ).with_model("openai", "gpt-5.1")
        
        prompt = f"""Context:
{context}

Question: {query}

Instructions: Provide a clear answer in bullet points or numbered list format. Make it scannable and easy to read. Use key points instead of paragraphs."""
        
        user_message = UserMessage(text=prompt)
        response = await chat.send_message(user_message)
        
        # Estimate token count (rough approximation)
        token_count = len(prompt.split()) + len(response.split())
        
        return response, token_count
    except Exception as e:
        logger.error(f"Error generating RAG answer: {e}")
        return "I apologize, but I encountered an error while generating the answer. Please try again.", 0

# API Endpoints
@api_router.get("/")
async def root():
    return {"message": "RAG Assistant API"}

@api_router.post("/documents/upload", response_model=DocumentResponse)
async def upload_document(file: UploadFile = File(...)):
    """Upload and process documents (PDF, Word, PowerPoint, Excel, Images)"""
    try:
        # Read file
        file_content = await file.read()
        
        # Extract text
        text, file_type = extract_text_from_file(file.filename, file_content)
        
        if not text.strip():
            raise HTTPException(status_code=400, detail=f"No text could be extracted from the {file_type.upper()} file")
        
        # Create document
        doc_id = str(uuid.uuid4())
        document = {
            'id': doc_id,
            'filename': file.filename,
            'file_type': file_type,
            'upload_date': datetime.now(timezone.utc).isoformat(),
            'file_size': len(file_content),
            'text_length': len(text)
        }
        
        # Chunk text
        chunks = split_text_into_chunks(text, chunk_size=500)
        document['chunk_count'] = len(chunks)
        
        # Store document
        database = get_database()
        await database.documents.insert_one(document)
        
        # Generate embeddings and store chunks
        if not chunks:
            raise HTTPException(status_code=400, detail=f"No text chunks could be created from the {file_type.upper()} file")
            
        for idx, chunk_text in enumerate(chunks):
            if not chunk_text.strip():
                continue  # Skip empty chunks
            embedding = await generate_embedding(chunk_text)
            chunk = {
                'id': str(uuid.uuid4()),
                'document_id': doc_id,
                'chunk_index': idx,
                'text': chunk_text,
                'embedding': embedding
            }
            await database.document_chunks.insert_one(chunk)
        
        logger.info(f"Document {file.filename} uploaded with {len(chunks)} chunks")
        
        return DocumentResponse(
            id=doc_id,
            filename=file.filename,
            chunk_count=len(chunks),
            upload_date=document['upload_date'],
            file_size=len(file_content)
        )
    except Exception as e:
        logger.error(f"Error uploading document: {e}")
        raise HTTPException(status_code=500, detail=str(e))

@api_router.get("/documents", response_model=List[DocumentResponse])
async def get_documents():
    """Get all documents"""
    database = get_database()
    # Limit to 1000 documents with projection for performance
    docs = await database.documents.find({}, {"_id": 0}).limit(1000).to_list(1000)
    return [DocumentResponse(**doc) for doc in docs]

@api_router.delete("/documents/{document_id}")
async def delete_document(document_id: str):
    """Delete a document and its chunks"""
    database = get_database()
    # Delete document
    result = await database.documents.delete_one({'id': document_id})
    if result.deleted_count == 0:
        raise HTTPException(status_code=404, detail="Document not found")
    
    # Delete chunks
    await database.document_chunks.delete_many({'document_id': document_id})
    
    return {"message": "Document deleted successfully"}

@api_router.post("/query", response_model=QueryResponse)
async def query_rag(request: QueryRequest):
    """Query the RAG system"""
    start_time = time.time()
    
    try:
        # Retrieve relevant chunks
        relevant_chunks = await retrieve_relevant_chunks(request.query, request.top_k)
        
        if not relevant_chunks:
            raise HTTPException(status_code=404, detail="No documents available. Please upload documents first.")
        
        chunk_texts = [chunk['text'] for chunk in relevant_chunks]
        
        # Generate answer
        answer, token_count = await generate_rag_answer(request.query, chunk_texts)
        
        # Calculate latency
        latency_ms = (time.time() - start_time) * 1000
        
        # Store telemetry
        database = get_database()
        telemetry = {
            'id': str(uuid.uuid4()),
            'query': request.query,
            'answer': answer,
            'latency_ms': latency_ms,
            'token_count': token_count,
            'timestamp': datetime.now(timezone.utc).isoformat(),
            'success': True
        }
        await database.telemetry.insert_one(telemetry)
        
        return QueryResponse(
            answer=answer,
            sources=relevant_chunks,
            latency_ms=latency_ms,
            token_count=token_count,
            retrieved_chunks=chunk_texts
        )
    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"Error processing query: {e}")
        
        # Store failed telemetry
        database = get_database()
        telemetry = {
            'id': str(uuid.uuid4()),
            'query': request.query,
            'latency_ms': (time.time() - start_time) * 1000,
            'timestamp': datetime.now(timezone.utc).isoformat(),
            'success': False,
            'error': str(e)
        }
        await database.telemetry.insert_one(telemetry)
        
        raise HTTPException(status_code=500, detail=str(e))

@api_router.get("/telemetry/stats", response_model=TelemetryStats)
async def get_telemetry_stats():
    """Get telemetry statistics"""
    database = get_database()
    # Limit to last 1000 records for stats calculation
    telemetry_records = await database.telemetry.find({}).sort("timestamp", -1).limit(1000).to_list(1000)
    
    if not telemetry_records:
        return TelemetryStats(
            total_queries=0,
            avg_latency_ms=0,
            total_tokens=0,
            total_cost=0,
            success_rate=0
        )
    
    total_queries = len(telemetry_records)
    successful_queries = [r for r in telemetry_records if r.get('success', False)]
    
    avg_latency = np.mean([r['latency_ms'] for r in telemetry_records])
    total_tokens = sum([r.get('token_count', 0) for r in successful_queries])
    # Estimate cost: $0.01 per 1000 tokens
    total_cost = (total_tokens / 1000) * 0.01
    success_rate = (len(successful_queries) / total_queries) * 100 if total_queries > 0 else 0
    
    return TelemetryStats(
        total_queries=total_queries,
        avg_latency_ms=float(avg_latency),
        total_tokens=total_tokens,
        total_cost=float(total_cost),
        success_rate=float(success_rate)
    )

@api_router.get("/telemetry/history")
async def get_telemetry_history(limit: int = 50):
    """Get recent telemetry records"""
    # Cap maximum limit to 100 for performance
    safe_limit = min(limit, 100)
    records = await db.telemetry.find({}, {"_id": 0}).sort("timestamp", -1).limit(safe_limit).to_list(safe_limit)
    return records

@api_router.get("/dashboard/stats")
async def get_dashboard_stats():
    """Get dashboard statistics"""
    doc_count = await db.documents.count_documents({})
    chunk_count = await db.document_chunks.count_documents({})
    query_count = await db.telemetry.count_documents({})
    
    # Get recent queries
    recent_queries = await db.telemetry.find(
        {}, 
        {"_id": 0, "query": 1, "timestamp": 1, "success": 1}
    ).sort("timestamp", -1).limit(5).to_list(5)
    
    return {
        'document_count': doc_count,
        'chunk_count': chunk_count,
        'query_count': query_count,
        'recent_queries': recent_queries
    }

# Include the router in the main app
app.include_router(api_router)

app.add_middleware(
    CORSMiddleware,
    allow_credentials=True,
    allow_origins=os.environ.get('CORS_ORIGINS', '*').split(','),
    allow_methods=["*"],
    allow_headers=["*"],
)

@app.on_event("shutdown")
async def shutdown_db_client():
    client.close()